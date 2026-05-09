import json
import re
import time
from pathlib import Path

from declutter_bot.core.utils import KEYWORD_TO_SUBJECT, SEED_MAP, SUBJECT_DESCRIPTIONS

GEMMA_MODEL = "gemma3:4b"             # local Ollama: text classification
GEMMA_VISION_MODEL = "gemma4:31b"    # local Ollama: vision (needs ~20 GB RAM)
GEMMA_GOOGLE_MODEL = "gemma-4-26b-a4b-it"  # Google AI API: text + vision (MoE, 4B active params)

# Image extensions that Gemma 4 Vision can process
VISUAL_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"}

# Folder names too generic to use as subject signals
GENERIC_FOLDER_NAMES = {
    "documents", "downloads", "desktop", "files", "school",
    "homework", "assignments", "work", "misc", "stuff",
    "new folder", "untitled", "folder", "archive", "backup"
}


# ---------------------------------------------------------
# Keyword extraction helper
# ---------------------------------------------------------

def extract_keywords(text: str) -> list[str]:
    """
    Split text into lowercase keyword tokens.
    e.g. "AP_Bio_Chapter3_Notes" → ["ap", "bio", "chapter", "notes"]
    """
    text = text.lower()
    text = re.sub(r'\.\w+$', '', text)
    tokens = re.split(r'[^a-z]+', text)
    return [t for t in tokens if len(t) > 1]


# ---------------------------------------------------------
# Scoring engines
# ---------------------------------------------------------

def score_text(text: str) -> dict[str, int]:
    """
    Score text against SEED_MAP.
    Multi-word phrase match → 2 points. Single-word → 1 point.
    Each unique keyword capped at 3 points to prevent repetition bias.
    """
    text_lower = text.lower()
    scores: dict[str, int] = {}

    for kw, subject in KEYWORD_TO_SUBJECT.items():
        weight = 2 if " " in kw else 1
        count = len(re.findall(r"\b" + re.escape(kw) + r"\b", text_lower))
        if count:
            scores[subject] = scores.get(subject, 0) + min(count * weight, 3)

    return scores


def confidence_from_scores(scores: dict[str, int]) -> tuple[str | None, float]:
    """
    Compute (winner_subject, confidence) from scores dict.
    confidence = winner_score / total_score.
    Ties → capped at 0.3 so the next group gets a chance.
    """
    if not scores:
        return None, 0.0
    total = sum(scores.values())
    winner = max(scores, key=scores.get)
    winner_score = scores[winner]
    sorted_vals = sorted(scores.values(), reverse=True)
    if len(sorted_vals) > 1 and sorted_vals[0] == sorted_vals[1]:
        return winner, 0.3
    return winner, winner_score / total


def _is_ambiguous(scores: dict[str, int]) -> bool:
    """Return True if the top two subjects are within 0.15 confidence of each other."""
    if len(scores) < 2:
        return False
    total = sum(scores.values())
    if total == 0:
        return True
    vals = sorted(scores.values(), reverse=True)
    return (vals[0] / total - vals[1] / total) < 0.15


# ---------------------------------------------------------
# File text readers
# ---------------------------------------------------------

def extract_text_from_bytes(data: bytes, ext: str, max_chars: int = 2000) -> str:
    """Extract plain text from binary file content (docx, pptx, or raw decode)."""
    try:
        if ext == ".docx":
            import docx, io
            doc = docx.Document(io.BytesIO(data))
            return " ".join(p.text for p in doc.paragraphs)[:max_chars]
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(data))
            parts = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            return " ".join(parts)[:max_chars]
        else:
            return data.decode("utf-8", errors="ignore")[:max_chars]
    except Exception:
        return ""


def _read_file_text(file_path: str, max_chars: int) -> str:
    """Read first max_chars of a local file as text."""
    ext = Path(file_path).suffix.lower()
    try:
        if ext in (".docx", ".pptx", ".ppt"):
            with open(file_path, "rb") as f:
                return extract_text_from_bytes(f.read(), ext, max_chars)
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(file_path, strict=False)
            parts = []
            total = 0
            for page in reader.pages:
                chunk = page.extract_text() or ""
                parts.append(chunk)
                total += len(chunk)
                if total >= max_chars:
                    break
            return " ".join(parts)[:max_chars]
        with open(file_path, "r", errors="ignore") as f:
            return f.read(max_chars)
    except (OSError, PermissionError, IsADirectoryError):
        return ""
    except Exception:
        return ""


def _normalize_name(name: str) -> str:
    """Replace underscores, hyphens, and digits with spaces so \b word boundaries work."""
    return re.sub(r'[_\-0-9]+', ' ', name)


def _metadata_text(file_path: str) -> str:
    """Build combined folder + filename text for Group A scoring."""
    path = Path(file_path)
    parts = [_normalize_name(path.stem)] + [
        _normalize_name(p.name) for p in path.parents
        if p.name and p.name.lower() not in GENERIC_FOLDER_NAMES
    ]
    return " ".join(parts)


# ---------------------------------------------------------
# Group A — metadata only (folder + filename keywords)
# ---------------------------------------------------------

def classify_group_a(file_path: str, display_name: str | None = None) -> tuple[str | None, float, dict[str, int]]:
    """
    Score folder names + filename against SEED_MAP.
    Pass display_name for Drive files where file_path is an opaque ID.
    Returns (subject, confidence, raw_scores).
    """
    text = _normalize_name(display_name) if display_name else _metadata_text(file_path)
    scores = score_text(text)
    subject, confidence = confidence_from_scores(scores)
    return subject, confidence, scores


# ---------------------------------------------------------
# Group B — content keyword scoring
# ---------------------------------------------------------

def classify_group_b(file_path: str, max_chars: int = 500) -> tuple[str | None, float, dict[str, int]]:
    """
    Read first max_chars of local file and score against SEED_MAP.
    Returns (subject, confidence, raw_scores).
    """
    content = _read_file_text(file_path, max_chars)
    if not content.strip():
        return None, 0.0, {}
    scores = score_text(content)
    subject, confidence = confidence_from_scores(scores)
    return subject, confidence, scores


# ---------------------------------------------------------
# Gemma inference helpers
# ---------------------------------------------------------

def _call_gemma(prompt: str) -> str | None:
    """
    Text classification via Gemma 3 on Ollama (local, private).
    Falls back to Gemma 4 via Google AI API if Ollama is unavailable.
    """
    try:
        import ollama
        response = ollama.chat(model=GEMMA_MODEL, messages=[{"role": "user", "content": prompt}])
        if isinstance(response, dict):
            m = response.get("message")
            raw = (m.get("content") if isinstance(m, dict) else None) or response.get("content")
        elif hasattr(response, "message"):
            raw = response.message.content
        else:
            raw = None
        if raw:
            return raw
    except Exception:
        pass  # fall through to Google AI API

    import os
    api_key = os.environ.get("GOOGLE_API_KEY")
    if not api_key:
        return None
    try:
        from google import genai
        client = genai.Client(api_key=api_key)
        backoff = 2
        deadline = time.monotonic() + 600
        while True:
            try:
                return client.models.generate_content(model=GEMMA_GOOGLE_MODEL, contents=prompt).text
            except Exception as e:
                is_rate_limit = any(t in str(e).lower() for t in ("429", "quota", "rate", "resource exhausted"))
                remaining = deadline - time.monotonic()
                if not is_rate_limit or remaining <= 0:
                    raise
                time.sleep(min(backoff, remaining))
                backoff = min(backoff * 2, 60)
    except Exception:
        return None


def _call_gemma_vision(
    prompt: str,
    image_bytes: bytes,
    image_mime_type: str = "image/jpeg",
) -> str | None:
    """
    Vision classification via Gemma 4 on Google AI API (multimodal).
    Gemma 3 does not support vision — this path always requires GOOGLE_API_KEY.
    Falls back to Ollama gemma4:31b only if available locally (~20 GB RAM needed).
    """
    import os
    api_key = os.environ.get("GOOGLE_API_KEY")
    if api_key:
        try:
            from google import genai
            from google.genai import types
            client = genai.Client(api_key=api_key)
            contents = [prompt, types.Part.from_bytes(data=image_bytes, mime_type=image_mime_type)]
            backoff = 2
            deadline = time.monotonic() + 600
            while True:
                try:
                    return client.models.generate_content(model=GEMMA_GOOGLE_MODEL, contents=contents).text
                except Exception as e:
                    is_rate_limit = any(t in str(e).lower() for t in ("429", "quota", "rate", "resource exhausted"))
                    remaining = deadline - time.monotonic()
                    if not is_rate_limit or remaining <= 0:
                        raise
                    time.sleep(min(backoff, remaining))
                    backoff = min(backoff * 2, 60)
        except Exception:
            pass  # fall through to local Ollama vision model

    try:
        import ollama
        msg: dict = {"role": "user", "content": prompt, "images": [image_bytes]}
        response = ollama.chat(model=GEMMA_VISION_MODEL, messages=[msg])
        if isinstance(response, dict):
            m = response.get("message")
            raw = (m.get("content") if isinstance(m, dict) else None) or response.get("content")
        elif hasattr(response, "message"):
            raw = response.message.content
        else:
            raw = None
        return raw
    except Exception:
        return None


def _parse_gemma_json(
    raw: str,
    candidates: list[str],
) -> tuple[str | None, float, str | None]:
    """Parse Gemma's JSON response and validate subject against the offered candidates."""
    valid_subjects = set(SEED_MAP.keys()) | {"other"}
    subject = None
    confidence = 0.0
    also_could_be = None

    try:
        json_match = re.search(r'\{.*\}', raw, re.DOTALL)
        if json_match:
            parsed = json.loads(json_match.group())
            subject = parsed.get("subject", "").lower().strip()
            confidence = float(parsed.get("confidence", 0))
            also_raw = parsed.get("also_could_be")
            if also_raw and str(also_raw).lower() not in ("null", "none", ""):
                also_could_be = str(also_raw).lower().strip()
    except Exception:
        m_sub = re.search(r'"subject"\s*:\s*"([^"]+)"', raw, re.I)
        m_conf = re.search(r'"confidence"\s*:\s*([0-9.]+)', raw, re.I)
        subject = m_sub.group(1).lower().strip() if m_sub else None
        confidence = float(m_conf.group(1)) if m_conf else 0.0

    allowed = set(candidates) | {"other"}
    if subject not in allowed:
        subject = None
    if also_could_be and also_could_be not in valid_subjects:
        also_could_be = None
    if also_could_be in (subject, "other", None):
        also_could_be = None

    return subject, confidence, also_could_be


# ---------------------------------------------------------
# Group C — Gemma (constrained contextual arbitrator)
#
# Called when Groups A+B confidence < 0.9 or result is ambiguous.
# Narrows to the top-2 keyword candidates and asks Gemma to reason
# step-by-step about the file's PURPOSE before committing to an answer.
# When confidence is below 0.75, returns also_could_be so the UI can
# flag the file for student review.
# ---------------------------------------------------------

def classify_group_c(
    text: str,
    file_path: str,
    keyword_scores: dict[str, int],
) -> tuple[str | None, float, str | None]:
    """
    Ask Gemma to arbitrate between the top-2 keyword candidates using chain-of-thought.

    Returns (subject, confidence, also_could_be).
    also_could_be is the runner-up subject when Gemma was uncertain (confidence < 0.75),
    otherwise None. Returns (None, 0.0, None) if Ollama is unavailable.
    """
    path = Path(file_path)
    filename = path.stem
    folder_parts = [
        p.name for p in path.parents
        if p.name and p.name.lower() not in GENERIC_FOLDER_NAMES
    ]
    folder = " > ".join(reversed(folder_parts[:3])) if folder_parts else "unknown"

    # Narrow to top-2 candidates — Gemma chooses between these, not all 12 subjects
    valid_subjects = set(SEED_MAP.keys()) | {"other"}
    ranked = sorted(keyword_scores.items(), key=lambda x: x[1], reverse=True)
    candidates = [s for s, _ in ranked if s in valid_subjects][:2]
    if not candidates:
        candidates = ["english", "history"]

    candidate_lines = []
    for i, subj in enumerate(candidates):
        label = chr(ord("A") + i)
        desc = SUBJECT_DESCRIPTIONS.get(subj, f"Files relating to {subj}.")
        candidate_lines.append(f'  {label}) "{subj}" — {desc}')
    other_label = chr(ord("A") + len(candidates))
    candidate_lines.append(f'  {other_label}) "other" — Does not fit any subject above.')
    candidates_block = "\n".join(candidate_lines)

    content_section = (
        f'\nFILE CONTENT (first 3000 characters):\n"""\n{text[:3000]}\n"""'
        if text.strip() else "\n(No readable text content in this file)"
    )

    prompt = f"""You are a school file classifier. A student created this file for one of their classes.

FILE
  Filename : "{filename}"
  Folder   : "{folder}"
{content_section}

CANDIDATES — you must choose exactly one:
{candidates_block}

STEP-BY-STEP REASONING
Work through the following before deciding. Write your reasoning for each step:

  Step 1 — Document type: What kind of document is this?
           (e.g. essay, study notes, problem set, lab report, schedule, reading, worksheet)

  Step 2 — Student task: What was the student doing when they made this file?
           Focus on the TASK, not the topic.
           Example: "writing an analytical essay" vs "taking notes on a historical event"

  Step 3 — Class match: Which candidate subject best matches that task?
           Remember — topic ≠ subject:
           • An essay ABOUT a historical event → "english" (the task is writing)
           • Notes ON a historical event → "history" (the task is history study)
           • A problem set using physics words → "math" (the task is solving equations)

  Step 4 — Decision:
           • If the file is a school ASSIGNMENT (essay, notes, worksheet, problem set,
             lab report, reading, debate, project) → pick the best-matching candidate.
           • If the file is administrative or non-academic (permission slip, roster,
             schedule, announcement, sign-up form, event flyer, recipe, personal letter)
             → pick "other".
           • If it is a close call between two academic candidates (confidence below 0.75),
             name the runner-up in "also_could_be" so the student can review.

After your reasoning, output ONLY this JSON on the final line (no markdown, no extra text):
{{"subject": "<exact name from candidates>", "confidence": <0.0-1.0>, "task_type": "<document type from step 1>", "also_could_be": "<runner-up subject or null>"}}"""

    raw = _call_gemma(prompt)
    if not raw:
        return None, 0.0, None
    return _parse_gemma_json(raw, candidates)


# ---------------------------------------------------------
# Group C (visual) — Gemma 4 Vision for image files
#
# Called for photos, screenshots, and scanned handwritten notes.
# Reads the image bytes and asks Gemma to describe the content,
# then classify into the top keyword candidates (or common subjects
# when no keyword signal is available).
# ---------------------------------------------------------

def classify_group_c_visual(
    file_path: str,
    keyword_scores: dict[str, int],
) -> tuple[str | None, float, str | None]:
    """
    Ask Gemma 4 Vision to classify an image (photo, screenshot, scanned notes).

    Returns (subject, confidence, also_could_be).
    Returns (None, 0.0, None) if Ollama is unavailable or the image can't be read.
    """
    path = Path(file_path)
    filename = path.stem

    valid_subjects = set(SEED_MAP.keys()) | {"other"}
    ranked = sorted(keyword_scores.items(), key=lambda x: x[1], reverse=True)
    candidates = [s for s, _ in ranked if s in valid_subjects][:2]
    if len(candidates) < 2:
        for subj in ["english", "math", "history", "science"]:
            if subj not in candidates:
                candidates.append(subj)
            if len(candidates) >= 4:
                break

    candidate_lines = []
    for i, subj in enumerate(candidates):
        label = chr(ord("A") + i)
        desc = SUBJECT_DESCRIPTIONS.get(subj, f"Files relating to {subj}.")
        candidate_lines.append(f'  {label}) "{subj}" — {desc}')
    other_label = chr(ord("A") + len(candidates))
    candidate_lines.append(f'  {other_label}) "other" — Personal photo or image with no academic content.')
    candidates_block = "\n".join(candidate_lines)

    prompt = f"""You are a school file classifier. A student saved this image file.

FILE: "{filename}"

Look at the image carefully. It may be a photo of handwritten notes, a whiteboard, a diagram, a worksheet, or a screenshot of schoolwork.

CANDIDATES — choose exactly one:
{candidates_block}

STEP-BY-STEP REASONING
  Step 1 — What do you see? Describe the image content briefly.
           (handwriting, equations, diagrams, text, graphs, labels, etc.)

  Step 2 — Is this school-related? What subject clues are visible?

  Step 3 — Which candidate subject best matches the visible content?

  Step 4 — Decision:
           • If the image shows school notes, assignments, diagrams, or academic content
             → pick the best-matching candidate.
           • If it's a personal photo (people, social events, scenery) with no academic content
             → pick "other".
           • If two subjects are equally plausible (confidence < 0.75), name the runner-up
             in "also_could_be".

After your reasoning, output ONLY this JSON on the final line (no markdown, no extra text):
{{"subject": "<exact name from candidates>", "confidence": <0.0-1.0>, "task_type": "<image type>", "also_could_be": "<runner-up subject or null>"}}"""

    try:
        import mimetypes
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type or not mime_type.startswith("image/"):
            mime_type = "image/jpeg"
        with open(file_path, "rb") as f:
            image_data = f.read()
    except (OSError, IOError):
        return None, 0.0, None

    raw = _call_gemma_vision(prompt, image_bytes=image_data, image_mime_type=mime_type)
    if not raw:
        return None, 0.0, None
    return _parse_gemma_json(raw, candidates)


# ---------------------------------------------------------
# Legacy pipeline — kept for tests and backward compatibility
# ---------------------------------------------------------

def classify_by_folder(file_path: str) -> str | None:
    path = Path(file_path)
    for folder in path.parents:
        name = folder.name.lower()
        if name in GENERIC_FOLDER_NAMES:
            continue
        for kw in extract_keywords(name):
            if kw in KEYWORD_TO_SUBJECT:
                return KEYWORD_TO_SUBJECT[kw]
    return None


def classify_by_filename(file_path: str) -> str | None:
    name = Path(file_path).stem
    for kw in extract_keywords(name):
        if kw in KEYWORD_TO_SUBJECT:
            return KEYWORD_TO_SUBJECT[kw]
    return None


def classify_by_seed_map(file_path: str) -> str | None:
    path = Path(file_path)
    all_text = " ".join([path.stem] + [p.name for p in path.parents]).lower()
    scores: dict[str, int] = {}
    for kw, subject in KEYWORD_TO_SUBJECT.items():
        if " " in kw:
            if re.search(r"\b" + re.escape(kw) + r"\b", all_text):
                scores[subject] = scores.get(subject, 0) + 2
    for token in extract_keywords(all_text):
        if token in KEYWORD_TO_SUBJECT:
            subject = KEYWORD_TO_SUBJECT[token]
            scores[subject] = scores.get(subject, 0) + 1
    return max(scores, key=scores.get) if scores else None


def classify_subject(file_path: str) -> str | None:
    return (
        classify_by_folder(file_path) or
        classify_by_filename(file_path) or
        classify_by_seed_map(file_path)
    )
